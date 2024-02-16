import os
import customtkinter 
from tkinter import filedialog, messagebox
from pptx import Presentation


#Extract text from pptx
def extract_text_from_pptx(file_path):
    try:
        presentation = Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {str(e)}")
        return None
    
#open file function
def open_file():
    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if file_path:
        extracted_text = extract_text_from_pptx(file_path)
        if extracted_text:
            save_text_to_file(extracted_text)

#folder function
def open_folder():
    folder_path = filedialog.askdirectory(title="Select Folder")
    if folder_path:
        for filename in os.listdir(folder_path):
            if filename.endswith(".pptx"):
                file_path = os.path.join(folder_path, filename)
                extracted_text = extract_text_from_pptx(file_path)
                if extracted_text:
                    save_text_to_file(extracted_text)
                    
#save text to files
def save_text_to_file(text):
    file_path = filedialog.asksaveasfilename(defaultextension=".txt", filetypes=[("Text files", "*.txt")])
    if file_path:
        try:
            with open(file_path, "w") as f:
                f.write(text)
            messagebox.showinfo("Success", "Text saved successfully!")
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred while saving the file: {str(e)}")

# GUI setup
app = customtkinter.CTk()
app.title("Powerpoint Text Extractor")
app.geometry("700x600")

open_button = customtkinter.CTkButton(app, height=50, width=100, text="Open .pptx File", command=open_file)
open_button.pack(padx=20, pady=20)

open_button2 = customtkinter.CTkButton(app, height=50, width=100, text="Open Folder", command=open_folder)
open_button2.pack(padx=30, pady=30)

#appearance
customtkinter.set_appearance_mode("Dark")

#warning label
label = customtkinter.CTkLabel(app, text='For the "Open Folder" option, ensure the folder only has powerpoint files or else errors will occur.', fg_color="transparent")
label.pack(padx=40, pady=40)

# loop
app.mainloop()
